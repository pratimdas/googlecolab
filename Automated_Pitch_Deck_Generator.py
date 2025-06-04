# ---
# jupyter:
#   jupytext:
#     text_representation:
#       extension: .py
#       format_name: percent
#       format_version: '1.3'
#       jupytext_version: 1.17.2
#   kernelspec:
#     display_name: Python 3
#     name: python3
# ---

# %% [markdown] id="view-in-github" colab_type="text"
# <a href="https://colab.research.google.com/github/pratimdas/googlecolab/blob/main/Automated_Pitch_Deck_Generator.ipynb" target="_parent"><img src="https://colab.research.google.com/assets/colab-badge.svg" alt="Open In Colab"/></a>

# %% id="Sp0l5Tdd3uzM"

# %% [markdown] id="KcJeUExH3xZ4"
# Here's a structured blueprint for building a simple, end-to-end PoC of the AI-powered pitch deck generator. Each step will be presented as a prompt to a code-generation LLM for implementation in a Google Colab environment, broken down into incremental chunks with clear context and progressive integration.Step-by-Step Blueprint for PoC Implementation
# Goal:
# Implement a simplified, harmonious end-to-end flow demonstrating core functionality:
#
# User input handling
# Basic LLM integration (simulated Gemini functionality)
# Automated simple research & insight generation
# Automated, basic pitch deck creation (using a lightweight library, e.g., Python-pptx)
# Minimal UI for interaction (Google Colab widgets/forms)
# We will avoid complex integrations or features unsuitable for Colab.

# %% id="qMq-pJaj4Bd8"

# %% [markdown] id="rvMxOSkS4bzz"
# Chunk 1: Setup and Initial Environment Preparation
# Prompt 1: Colab Environment Setup

# %% colab={"base_uri": "https://localhost:8080/"} id="pT5fCDF24d9N" outputId="6b1d22e2-c206-4e68-efc3-9d74cd92659f"
# Install necessary libraries
# !pip install python-pptx requests pandas numpy

# Import required libraries
from pptx import Presentation
from pptx.util import Inches
import requests
import pandas as pd
import numpy as np

# Set up initial variables/placeholders for later use
user_input = None           # Placeholder for the user's input (e.g., pitch details)
research_data = None        # Placeholder for automated research & insight generation
pitch_deck = None           # Placeholder for the generated pitch deck (PowerPoint presentation)

print("Environment setup complete. Libraries installed and initial variables initialized.")


# %% [markdown] id="INL94CsO4z6h"
# Chunk 2: Simple User Input Collection
# Prompt 2: Collect User Inputs via Widgets

# %% colab={"base_uri": "https://localhost:8080/", "height": 195, "referenced_widgets": ["bb46517a707a484db356df8c496f4cd5", "18b2023a79cd4b29bf8e262ee3df0c6f", "9c0d5fe37ee941c0ae102db8acf7c1e5", "42268bbbcffe44f68886b5ee5c1c519d", "1d90e5d1752b4309beb792a9ca0f1c65", "9da17b9df1a1457c9a8fed786c8e13bc", "eb764b8381c640a88868d9d54cc26cff", "9649d01429654460936cd8c7aebd2ac3", "d0a394b51305416b8564ab519564277b", "1ad5f076b73f4b0eb72d61eb13a53e19", "a5fdbc088d12492ebf147c0f3f327095", "9b1e6735f07c48abb94eea565f77d05a", "622e9e242f8a469c9c5b3b1c59176670", "a35fe8b3c9d147909e960a2e8b15cec5"]} id="Awk9ZhvA5fMI" outputId="eaa2d962-01d2-46a0-c29a-d0633643f011"
import ipywidgets as widgets
from IPython.display import display

# Create widget for company name input (text)
company_name_widget = widgets.Text(
    value='',
    placeholder='Enter your company name',
    description='Company Name:',
    disabled=False
)

# Create widget for industry/domain focus input (dropdown)
industry_options = ['Technology', 'Healthcare', 'Finance', 'Retail', 'Other']
industry_widget = widgets.Dropdown(
    options=industry_options,
    value=industry_options[0],
    description='Industry:',
    disabled=False,
)

# Create widget for logo/branding asset upload (single file, images only)
logo_upload_widget = widgets.FileUpload(
    accept='image/*',  # Accept images only
    multiple=False,    # Single file upload
    description='Upload Logo'
)

# Display the widgets in the notebook
display(company_name_widget)
display(industry_widget)
display(logo_upload_widget)

# Create a submit button to capture the inputs
submit_button = widgets.Button(
    description='Submit',
    button_style='success'
)

output = widgets.Output()

def on_submit_clicked(b):
    with output:
        output.clear_output()
        # Store widget values in clearly named Python variables
        global company_name, industry_focus, logo_data
        company_name = company_name_widget.value
        industry_focus = industry_widget.value

        # For the file uploader: check if a file was uploaded and retrieve its content
        if logo_upload_widget.value:
            uploaded_file = list(logo_upload_widget.value.values())[0]
            logo_data = uploaded_file['content']
        else:
            logo_data = None

        print("Company Name:", company_name)
        print("Industry/Domain Focus:", industry_focus)
        print("Logo Uploaded:", "Yes" if logo_data else "No")

# Bind the click event to the submit button
submit_button.on_click(on_submit_clicked)

# Display the button and output area
display(submit_button, output)


# %% id="jLfEEGMf5udh"

# %% [markdown] id="UnJnj5_mJDD6"
# Chunk 3: Basic LLM Integration (Simulated Gemini)
# # Create a simple placeholder function to simulate Gemini LLM:
# - Accept the company name and industry/domain inputs
# - Return a basic mocked response containing: Company Overview, Business Objectives, Pain Points, Competitor Analysis, Market Opportunities, and Key Executives (dummy data is acceptable here).

# %% colab={"base_uri": "https://localhost:8080/", "height": 575} id="eNfYk-hhKMcv" outputId="36ec3e0f-1bd6-4dda-f5d5-af7f0dcd31be"
# !pip install --upgrade google-generativeai #Upgrade the google-generativeai package to the latest version

import os
import json
# Import the Python SDK
import google.generativeai as genai
from google.colab import userdata
global research_text;
# Configure the Gemini client with the API key stored as GOOGLE_API_KEY in Colab
api_key = userdata.get('GOOGLE_API_KEY')
if not api_key:
    raise ValueError("GOOGLE_API_KEY not found. Please set the secret in your Colab notebook.")
genai.configure(api_key=api_key)

def call_gemini(company_name, industry_focus):
    """
    Uses the Gemini SDK (via google-generativeai) to perform deep research on the specified
    company and industry, then generates a concise summary of the research findings.

    Parameters:
        company_name (str): The name of the company.
        industry_focus (str): The industry or domain focus.

    Returns:
        tuple: (research_text, summary_text)
               - research_text: Detailed insights generated by Gemini, including Company Overview,
                 Business Objectives, Pain Points, Competitor Analysis, Market Opportunities, and Key Executives.
               - summary_text: A concise overview summarizing the research findings.
    """
    # --- Step 1: Perform Deep Research ---
    research_prompt = (
        f"Perform deep research on '{company_name}' operating in the '{industry_focus}' industry. "
        "Provide detailed insights including Company Overview, Business Objectives, Pain Points, "
        "Competitor Analysis, Market Opportunities, and Key Executives."
    )
    model = genai.GenerativeModel()
    research_response = model.generate_content(research_prompt);


    if not research_response.text:
        raise ValueError("Failed to generate research text from Gemini.")

    research_text = research_response.text

    # --- Step 2: Generate a Concise Summary ---
    summary_prompt = (
        "Summarize the following research findings into a concise overview:\n\n" + research_text
    )

    summary_response = model.generate_content(summary_prompt);

    if not summary_response.text:
        raise ValueError("Failed to generate summary text from Gemini.")

    summary_text = summary_response.text

    # Save the results as global variables for later use in the notebook
    global deep_research_results, deep_research_summary
    deep_research_results = research_text
    deep_research_summary = summary_text

    # Display the generated summary using Gemini itself
    print("Deep Research Summary:")
    print(summary_text)

    return research_text, summary_text

# Example usage (uncomment and run after providing company_name and industry_focus from user inputs):
research_data, summary = call_gemini(company_name, industry_focus)

# %% [markdown] id="tNBGadXZb9Pt"
# 4: Insight Generation and Data Structuring
#
# # Parse the LLM's mocked output:
# - Store each section (Company Overview, Objectives, Pain Points, Competitor Analysis, etc.) in a structured dictionary or DataFrame.
# - Display the structured data clearly for review within Colab.

# %% [markdown] id="JrbxI2XIKF5J"
#

# %% colab={"base_uri": "https://localhost:8080/", "height": 1000} id="8O3ebDwzcf4a" outputId="057ff1ea-b326-4ff1-eabe-896e5acf561a"
import re
import pandas as pd

def parse_llm_output(text):
    """
    Parses the Gemini LLM output text to extract sections and their content.
    Expected format in the text:
      Section Title: Content
      (Each section is separated by a newline and ends when a new section begins.)

    Parameters:
        text (str): The full text output from the LLM containing research insights.

    Returns:
        dict: A dictionary with section titles as keys and their corresponding content as values.
    """
    # This regex looks for a section name ending with a colon, and captures everything until the next section starts
    pattern = re.compile(
        r'(?P<section>[A-Za-z ]+):\s*(?P<content>.*?)(?=\n[A-Za-z ]+:|$)',
        re.DOTALL
    )

    sections = pattern.finditer(text)
    parsed_data = {}
    for match in sections:
        section = match.group('section').strip()
        content = match.group('content').strip()
        parsed_data[section] = content
    return parsed_data

# Assuming 'deep_research_results' variable holds the Gemini research text from the previous step.
# Check if deep_research_results exists instead of research_text
if 'deep_research_results' not in globals():
    raise ValueError("deep_research_results not found. Please run the Gemini call cell first.")

# Parse the Gemini output into structured sections, using deep_research_results
structured_data = parse_llm_output(deep_research_results)

# Display the structured data in a clear textual format
print("Structured Research Data:\n")
for section, content in structured_data.items():
    print(f"{section}:")
    print(content)
    print("-" * 50)

# Also display the structured data in a DataFrame for easier review in Colab
df = pd.DataFrame(list(structured_data.items()), columns=["Section", "Content"])
display(df)

# %% [markdown] id="VxV_EH7wenwb"
# 5: Allow User Review and Simple Editing of Insights
#
# # Provide an interactive cell allowing users to:
# - View each insight section clearly (displayed via markdown)
# - Edit any section using text boxes
# - Update the stored insights with edited data upon confirmation

# %% colab={"base_uri": "https://localhost:8080/", "height": 219, "referenced_widgets": ["776ad753d8e547afb9ff3edb7e25459a", "3f0196feb0374981ba0944370fcd1bbf", "8ce215ffc0dc4a93af05eb2b4d03ccdb", "f61cb39e00d443ffb254dca878283cf4", "24872e1c582146948548eb9f7362dd4f", "f2ee6402dbce4f7c90fc5afced32c2df", "20e438c06c7a4800b4152744368160b2", "4510fe6efc45401b84d266a5e9c3bd0b", "35432d42d80c4593b65bf673e8a0a2c7", "0e8da01d7d9a4c14ae95519a99c883af", "1e854b5906094ae2b4e354c5691b94bc", "56e91b73d9db442bb74fcf66386e72a6", "01354866f30b4080b45cc8af74f38e20"]} id="fCub6RNeerUK" outputId="76b21a65-35e9-4241-edf2-1b932b0da51b"
import ipywidgets as widgets
from IPython.display import display

# Ensure that structured_data is available from the previous cell.
if 'structured_data' not in globals():
    raise ValueError("structured_data not found. Please run the previous cells (Prompt 4) first.")

# Create a dictionary to hold Textarea widgets for each insight section.
section_widgets = {}

# Create a list to hold the widget elements for layout.
widget_elements = []

# Iterate over each section in the structured_data dictionary.
for section, content in structured_data.items():
    # Instead of using IPython.display.Markdown (which is not a widget),
    # we use ipywidgets.HTML to display section titles.
    title_html = widgets.HTML(value=f"<h3>{section}</h3>")

    # Create a Textarea widget pre-filled with the current content.
    text_area = widgets.Textarea(
        value=content,
        placeholder=f"Edit {section} content here...",
        layout=widgets.Layout(width='100%', height='100px')
    )

    # Store the widget for later retrieval.
    section_widgets[section] = text_area

    # Append the title and text area to our widget elements list.
    widget_elements.append(title_html)
    widget_elements.append(text_area)

# Create an update button for confirming edits.
update_button = widgets.Button(
    description="Update Insights",
    button_style="success",
    tooltip="Click to update the insights with your edits."
)

# Output widget to display messages upon updating.
output_update = widgets.Output()

def update_insights(b):
    with output_update:
        output_update.clear_output()
        # Update the global structured_data with the new values from the text areas.
        for section, widget_instance in section_widgets.items():
            structured_data[section] = widget_instance.value

        print("Insights updated successfully!\n")
        # Optionally, display the updated sections.
        for section, content in structured_data.items():
            print(f"{section}:\n{content}")
            print("-" * 50)

# Bind the button click event to the update_insights function.
update_button.on_click(update_insights)

# Display all widgets together: the editing interface, update button, and output area.
display(widgets.VBox(widget_elements))
display(update_button)
display(output_update)



# %% [markdown] id="L4ruVs3pgxP7"
# 6: Basic Automated Pitch Deck Generation
#
# # Using python-pptx, automatically generate a simple pitch deck with slides for:
# 1. Cover Page (Company Name and Logo)
# 2. Executive Summary
# 3. Company Overview
# 4. Pain Points
# 5. Competitor Analysis
# 6. Market Opportunities
# 7. Proposed Strategy (basic placeholder)
# 8. Next Steps (basic placeholder)
#
# Populate slides directly from the reviewed and approved insights.

# %% colab={"base_uri": "https://localhost:8080/"} id="ll5Qvlyzg1C5" outputId="34584856-6a6f-4980-f6de-b484d9c0116d"
# !pip install --upgrade google-auth google-auth-oauthlib google-auth-httplib2
# Authenticate with Google APIs (Slides and Drive)
from google.colab import auth
auth.authenticate_user()

from googleapiclient.discovery import build

# Build the Slides API service
slides_service = build('slides', 'v1')

# If a logo was uploaded, use PyDrive to upload it and get a public URL.
logo_url = None
if 'logo_data' in globals() and logo_data:
    # !pip install -U -q PyDrive
    from pydrive.auth import GoogleAuth
    from pydrive.drive import GoogleDrive
    from oauth2client.client import GoogleCredentials

    gauth = GoogleAuth()
    gauth.credentials = GoogleCredentials.get_application_default()
    drive = GoogleDrive(gauth)

    logo_file = drive.CreateFile({'title': 'logo_temp.png'})
    logo_file.SetContentBinary(logo_data)
    logo_file.Upload()
    # Make the file public.
    logo_file.InsertPermission({
        'type': 'anyone',
        'value': 'anyone',
        'role': 'reader'
    })
    logo_url = f"https://drive.google.com/uc?id={logo_file['id']}"

# Create a new Google Slides presentation
presentation = slides_service.presentations().create(body={
    'title': f'Pitch Deck: {company_name}'
}).execute()
presentation_id = presentation.get('presentationId')
print("Created presentation with ID:", presentation_id)

# We'll build a list of batchUpdate requests to add slides and their content.
requests = []

# Utility function to add a slide with a title and content text box.
def add_slide(slide_id, title_text, content_text):
    # Create a blank slide.
    requests.append({
        "createSlide": {
            "objectId": slide_id,
            "insertionIndex": "0",
            "slideLayoutReference": {"predefinedLayout": "BLANK"}
        }
    })
    # Create a text box for the title.
    title_box_id = slide_id + "_title"
    requests.append({
        "createShape": {
            "objectId": title_box_id,
            "shapeType": "TEXT_BOX",
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {"height": {"magnitude": 50, "unit": "PT"},
                         "width": {"magnitude": 400, "unit": "PT"}},
                "transform": {"scaleX": 1, "scaleY": 1,
                              "translateX": 50, "translateY": 50, "unit": "PT"}
            }
        }
    })
    requests.append({
        "insertText": {
            "objectId": title_box_id,
            "insertionIndex": 0,
            "text": title_text
        }
    })
    # Create a text box for the content.
    content_box_id = slide_id + "_content"
    requests.append({
        "createShape": {
            "objectId": content_box_id,
            "shapeType": "TEXT_BOX",
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {"height": {"magnitude": 300, "unit": "PT"},
                         "width": {"magnitude": 400, "unit": "PT"}},
                "transform": {"scaleX": 1, "scaleY": 1,
                              "translateX": 50, "translateY": 120, "unit": "PT"}
            }
        }
    })
    requests.append({
        "insertText": {
            "objectId": content_box_id,
            "insertionIndex": 0,
            "text": content_text
        }
    })

# --- Slide 1: Cover Page ---
slide1_id = "slide1"
requests.append({
    "createSlide": {
        "objectId": slide1_id,
        "insertionIndex": "0",
        "slideLayoutReference": {"predefinedLayout": "BLANK"}
    }
})
# Add Company Name as the title.
cover_title_id = "slide1_title"
requests.append({
    "createShape": {
        "objectId": cover_title_id,
        "shapeType": "TEXT_BOX",
        "elementProperties": {
            "pageObjectId": slide1_id,
            "size": {"height": {"magnitude": 50, "unit": "PT"},
                     "width": {"magnitude": 400, "unit": "PT"}},
            "transform": {"scaleX": 1, "scaleY": 1,
                          "translateX": 50, "translateY": 50, "unit": "PT"}
        }
    }
})
requests.append({
    "insertText": {
        "objectId": cover_title_id,
        "insertionIndex": 0,
        "text": company_name
    }
})
# Add Industry/Domain as a subtitle.
cover_subtitle_id = "slide1_subtitle"
requests.append({
    "createShape": {
        "objectId": cover_subtitle_id,
        "shapeType": "TEXT_BOX",
        "elementProperties": {
            "pageObjectId": slide1_id,
            "size": {"height": {"magnitude": 50, "unit": "PT"},
                     "width": {"magnitude": 400, "unit": "PT"}},
            "transform": {"scaleX": 1, "scaleY": 1,
                          "translateX": 50, "translateY": 120, "unit": "PT"}
        }
    }
})
requests.append({
    "insertText": {
        "objectId": cover_subtitle_id,
        "insertionIndex": 0,
        "text": industry_focus
    }
})
# Insert logo image if available.
if logo_url:
    logo_object_id = "slide1_logo"
    requests.append({
        "createImage": {
            "objectId": logo_object_id,
            "url": logo_url,
            "elementProperties": {
                "pageObjectId": slide1_id,
                "size": {"height": {"magnitude": 100, "unit": "PT"},
                         "width": {"magnitude": 100, "unit": "PT"}},
                "transform": {"scaleX": 1, "scaleY": 1,
                              "translateX": 500, "translateY": 50, "unit": "PT"}
            }
        }
    })

# --- Slide 2: Executive Summary ---
add_slide("slide2", "Executive Summary", deep_research_summary)

# --- Slide 3: Company Overview ---
# Use the parsed content from structured_data.
company_overview = structured_data.get("Company Overview", "Not available")
add_slide("slide3", "Company Overview", company_overview)

# --- Slide 4: Pain Points ---
pain_points = structured_data.get("Pain Points", "Not available")
add_slide("slide4", "Pain Points", pain_points)

# --- Slide 5: Competitor Analysis ---
competitor_analysis = structured_data.get("Competitor Analysis", "Not available")
add_slide("slide5", "Competitor Analysis", competitor_analysis)

# --- Slide 6: Market Opportunities ---
market_opportunities = structured_data.get("Market Opportunities", "Not available")
add_slide("slide6", "Market Opportunities", market_opportunities)

# --- Slide 7: Proposed Strategy (Placeholder) ---
add_slide("slide7", "Proposed Strategy", "Placeholder: Proposed strategy details go here.")

# --- Slide 8: Next Steps (Placeholder) ---
add_slide("slide8", "Next Steps", "Placeholder: Next steps details go here.")

# Execute the batchUpdate request to add all slides and their elements.
body = {'requests': requests}
response = slides_service.presentations().batchUpdate(
    presentationId=presentation_id, body=body).execute()
print("Google Slides pitch deck created successfully!")

# Provide a link to view the presentation.
print("View your presentation: https://docs.google.com/presentation/d/{}/edit".format(presentation_id))

